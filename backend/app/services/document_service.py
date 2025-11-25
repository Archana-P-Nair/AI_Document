from docx import Document
from docx.shared import Pt, Inches
from pptx import Presentation
from pptx.util import Inches, Pt as PptxPt
from io import BytesIO
import os

def create_word_document(project, sections):
    """
    Create a Word document from project and sections
    """
    doc = Document()
    
    # Add title
    title = doc.add_heading(project.title, 0)
    title.alignment = 1  # Center alignment
    
    # Add topic/description
    doc.add_paragraph(f"Topic: {project.topic}")
    doc.add_paragraph()  # Empty line
    
    # Add each section
    for section in sorted(sections, key=lambda x: x.order):
        # Section heading
        doc.add_heading(section.title, 1)
        
        # Section content
        if section.content:
            doc.add_paragraph(section.content)
        else:
            doc.add_paragraph("[Content not generated yet]")
        
        # Add space between sections
        doc.add_paragraph()
    
    # Save to BytesIO
    file_stream = BytesIO()
    doc.save(file_stream)
    file_stream.seek(0)
    
    return file_stream

def create_powerpoint_presentation(project, sections):
    """
    Create a PowerPoint presentation from project and sections
    """
    prs = Presentation()
    
    # Set slide size (16:9)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = project.title
    subtitle.text = project.topic
    
    # Add content slides
    for section in sorted(sections, key=lambda x: x.order):
        # Use title and content layout
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        
        # Add title
        title = slide.shapes.title
        title.text = section.title
        
        # Add content
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()  # Clear default text
        
        if section.content:
            # Split content by lines and add as bullet points
            lines = section.content.strip().split('\n')
            for i, line in enumerate(lines):
                line = line.strip()
                if line:
                    # Remove bullet characters if present
                    line = line.lstrip('•-*').strip()
                    
                    if i == 0:
                        # First line
                        p = text_frame.paragraphs[0]
                        p.text = line
                        p.level = 0
                    else:
                        # Subsequent lines
                        p = text_frame.add_paragraph()
                        p.text = line
                        p.level = 0
        else:
            p = text_frame.paragraphs[0]
            p.text = "[Content not generated yet]"
    
    # Save to BytesIO
    file_stream = BytesIO()
    prs.save(file_stream)
    file_stream.seek(0)
    
    return file_stream